"""Local RAG knowledge base over YOUR documents.

Indexes the user's DOCUMENTS (md, txt, rst, html, pdf->text) of a project and, for a question,
retrieves the TOP-K most relevant fragments (with their SOURCE) so the agent can answer BASED on
them and CITE them. Light context: the model never receives whole documents, only the top-k
fragments.

100% local and free: LEXICAL TF-IDF retrieval in pure Python (no model downloads, robust).
Design note: TF-IDF (not neural embeddings) was chosen for robustness and to avoid depending on a
fragile model download; the design is 'embedding-ready' (just swap the scorer). PDFs via `pdftotext`
(poppler) with robust handling of malformed PDFs. The base lives INSIDE the project (`.knowledge/`),
so each project has its own.
"""
import json
import math
import os
import re
import subprocess  # nosec B404: for `pdftotext`/`pandoc` (extract text from docs): fixed list, no shell
from collections import Counter
from pathlib import Path

KB_DIR = ".knowledge"
INDEX = "index.json"
DOC_EXTS = {".md", ".txt", ".rst", ".markdown", ".text", ".html", ".htm", ".pdf", ".csv", ".log",
            ".docx", ".xlsx", ".pptx"}   # +office docs (docx via pandoc; xlsx/pptx if lib present)
MAX_BYTES_DOC = 5_000_000      # don't read huge documents
CHUNK_CHARS = 1000             # fragment size
CHUNK_OVERLAP = 150            # overlap between fragments (don't lose context at the edges)
MAX_CHUNKS = 4000              # cap on indexed fragments (don't grow without bound)
PDF_TIMEOUT = 20               # cap for pdftotext (a malformed PDF won't hang the index)
PANDOC_TIMEOUT = 25            # cap for pandoc (a malformed docx won't hang the index)
_EXCLUDE_DIRS = {KB_DIR, ".snapshots", ".git", ".hg", ".svn", "__pycache__", ".pytest_cache",
                 "node_modules", ".venv", "venv"}

_RX_TOKEN = re.compile(r"[a-z0-9]+", re.IGNORECASE)


def _tokenize(text):
    return [t for t in _RX_TOKEN.findall((text or "").lower()) if len(t) >= 2]


# --- Text extraction by type (robust) -----------------------------------------
def _pdf_text(p: Path):
    try:
        r = subprocess.run(["pdftotext", "-q", str(p), "-"], capture_output=True,  # noqa: S603,S607: fixed command (no shell); path is a project doc; pdftotext via PATH on purpose  # nosec B603 B607
                            timeout=PDF_TIMEOUT)
        # bound the PDF text like the rest (a PDF with tons of text won't bloat the index)
        return r.stdout[:MAX_BYTES_DOC].decode("utf-8", errors="replace") if r.returncode == 0 else ""
    except (OSError, subprocess.TimeoutExpired):
        return ""              # no poppler or malformed PDF -> skip (doesn't break the index)


def _html_text(text):
    try:
        import trafilatura
        return trafilatura.extract(text) or re.sub(r"<[^>]+>", " ", text)
    except Exception:  # noqa: BLE001
        return re.sub(r"<[^>]+>", " ", text)


def _docx_text(p: Path):
    """`.docx` -> text via pandoc (already installed). Degrades honestly ('') if pandoc is missing or
    the file is malformed (never breaks the index)."""
    try:
        r = subprocess.run(["pandoc", "-f", "docx", "-t", "plain", "--wrap=none", str(p)],  # noqa: S603,S607: fixed command (no shell); path is a project doc; pandoc via PATH on purpose  # nosec B603 B607
                           capture_output=True, timeout=PANDOC_TIMEOUT)
        return r.stdout[:MAX_BYTES_DOC].decode("utf-8", errors="replace") if r.returncode == 0 else ""
    except (OSError, subprocess.TimeoutExpired):
        return ""


def _xlsx_text(p: Path):
    """`.xlsx` -> text via openpyxl IF installed; otherwise degrades honestly (''): the sheet is
    skipped without breaking anything (install `openpyxl` to index spreadsheets)."""
    try:
        import openpyxl
    except ImportError:
        return ""
    try:
        wb = openpyxl.load_workbook(p, read_only=True, data_only=True)
    except Exception:  # noqa: BLE001: malformed xlsx -> skip
        return ""
    parts = []
    try:
        for sheet in wb.worksheets:
            parts.append(f"# {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    parts.append("\t".join(cells))
                if sum(len(x) for x in parts) > MAX_BYTES_DOC:
                    break
    except Exception:  # noqa: BLE001: sheet corrupted mid-read -> return what was read, don't crash
        pass
    finally:
        try:
            wb.close()
        except Exception:  # noqa: BLE001
            pass
    return "\n".join(parts)[:MAX_BYTES_DOC]


def _pptx_text(p: Path):
    """`.pptx` -> text via python-pptx IF installed; otherwise degrades honestly ('')."""
    try:
        from pptx import Presentation
    except ImportError:
        return ""
    try:
        pres = Presentation(str(p))
    except Exception:  # noqa: BLE001
        return ""
    parts = []
    try:
        for i, slide in enumerate(pres.slides, 1):
            parts.append(f"# Slide {i}")
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    parts.append(shape.text_frame.text)
    except Exception:  # noqa: BLE001: slide/shape corrupted mid-read -> return what was read, don't crash
        pass
    return "\n".join(parts)[:MAX_BYTES_DOC]


def _extract_text(p: Path):
    ext = p.suffix.lower()
    try:
        if ext == ".pdf":
            return _pdf_text(p)
        if ext == ".docx":
            return _docx_text(p)
        if ext == ".xlsx":
            return _xlsx_text(p)
        if ext == ".pptx":
            return _pptx_text(p)
        raw = p.read_bytes()
        if len(raw) > MAX_BYTES_DOC:
            raw = raw[:MAX_BYTES_DOC]
        text = raw.decode("utf-8", errors="replace")
        if ext in (".html", ".htm"):
            return _html_text(text)
        return text
    except OSError:
        return ""


def _documents(base: Path):
    """RELATIVE paths of the user's DOCUMENTS under `base` (by extension), pruning system/agent dirs
    and never leaving `base` (not even through a symlink pointing OUTSIDE the project, which avoids
    leaking external content into the model's context)."""
    base = Path(base)
    base_real = base.resolve()
    for root, dirs, files in os.walk(base, followlinks=False):
        dirs[:] = [d for d in sorted(dirs) if d not in _EXCLUDE_DIRS]
        for name in sorted(files):
            p = Path(root) / name
            if p.suffix.lower() not in DOC_EXTS or name.startswith("."):
                continue
            try:
                real = p.resolve()
                if base_real not in real.parents and real != base_real:   # escaping symlink: don't index
                    continue
            except OSError:
                continue
            yield p.relative_to(base)


def _chunk(text, source):
    """Split the text into overlapping fragments, annotating the source and approximate line. Counts
    lines INCREMENTALLY (it used to be O(n^2) due to `count('\\n',0,i)` on each iteration)."""
    text = re.sub(r"[ \t]+", " ", text or "")
    if not text.strip():
        return []
    chunks = []
    i, n, prev, line = 0, len(text), 0, 1
    step = max(1, CHUNK_CHARS - CHUNK_OVERLAP)
    while i < n and len(chunks) < MAX_CHUNKS:
        line += text.count("\n", prev, i)      # only the NEW span (incremental -> O(n) total)
        prev = i
        chunks.append({"source": str(source), "line": line, "text": text[i:i + CHUNK_CHARS].strip()})
        i += step
    return chunks


# --- Index (persistent, per project) ------------------------------------------
def _index_path(base: Path):
    d = Path(base) / KB_DIR
    d.mkdir(parents=True, exist_ok=True)
    return d / INDEX


def _mtimes(base: Path):
    out = {}
    for rel in _documents(base):
        try:
            out[str(rel)] = (base / rel).stat().st_mtime
        except OSError:
            continue
    return out


def index_docs(base, force=False, cfg=None):
    """(Re)build the index if missing or the documents changed. Returns (n_docs, n_chunks).
    Optional: if EMBEDDINGS is enabled and the model is available, attach one embedding per chunk
    (semantic layer); with OFF/degraded, the index is IDENTICAL to the usual one (no `emb`)."""
    import embeddings
    base = Path(base)
    path = _index_path(base)
    current = _mtimes(base)
    want_emb = embeddings.enabled(cfg) and embeddings.available(cfg)
    emb_model = embeddings.model(cfg) if want_emb else None
    if not force and path.exists():
        try:
            prev = json.loads(path.read_text(encoding="utf-8"))
            docs_ok = prev.get("docs") == current
            emb_ok = (not want_emb) or (prev.get("emb_model") == emb_model)   # ON: reindex if missing
            if docs_ok and emb_ok:
                return len(current), len(prev.get("chunks", []))   # fresh index
        except (OSError, json.JSONDecodeError):
            pass
    chunks = []
    for rel in current:
        try:                                          # ISOLATE each doc: a problematic one (corrupt
            fragments = _chunk(_extract_text(base / rel), rel)   # docx/xlsx/pptx, misbehaving lib...)
        except Exception:  # noqa: BLE001,S112  # nosec B112                  # is SKIPPED, not fatal to the index.
            continue
        for ch in fragments:
            if ch["text"]:
                ch["tf"] = dict(Counter(_tokenize(ch["text"])))
                chunks.append(ch)
                if len(chunks) >= MAX_CHUNKS:
                    break
        if len(chunks) >= MAX_CHUNKS:
            break
    df = Counter()
    for ch in chunks:
        df.update(ch["tf"].keys())
    n = max(1, len(chunks))
    idf = {t: math.log(1 + n / c) for t, c in df.items()}
    # L2 norm of each fragment in the tf-idf space (for the cosine)
    for ch in chunks:
        ch["norm"] = math.sqrt(sum((f * idf.get(t, 0.0)) ** 2 for t, f in ch["tf"].items())) or 1.0
    index = {"docs": current, "idf": idf, "chunks": chunks, "n": len(chunks)}
    if want_emb and chunks:                          # attach one embedding per chunk (semantic layer)
        vecs = embeddings.embed_passages([ch["text"] for ch in chunks], cfg)
        if vecs and len(vecs) == len(chunks):
            for ch, v in zip(chunks, vecs):
                ch["emb"] = v
            index["emb_model"] = emb_model           # mark that the index carries embeddings (freshness)
    tmp = path.with_suffix(".json.tmp")
    tmp.write_text(json.dumps(index, ensure_ascii=False), encoding="utf-8")
    os.replace(tmp, path)
    return len(current), len(chunks)


def _load_index(base: Path):
    path = _index_path(base)
    if not path.exists():
        return None
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return None


def has_index(base):
    """Does the project have indexable DOCUMENTS? (to decide whether a question 'about my documents'
    should be answered from the knowledge base)."""
    try:
        return any(True for _ in _documents(Path(base)))
    except OSError:
        return False


def indexed_names(base):
    """File names (relative) of the project's DOCUMENTS (to detect whether a question mentions a file
    that exists: 'according to my CV' with CV_....pdf indexed)."""
    try:
        return [str(rel) for rel in _documents(Path(base))]
    except OSError:
        return []


def _hybrid_search(index, query, only_source, k):
    """Combines the LEXICAL cosine (TF-IDF) with the SEMANTIC one (embeddings, rescaled over their
    baseline), which also retrieves across languages. Same light context: returns the top-k with a
    citation. None if the query could not be embedded (the caller falls back to TF-IDF, honest
    degradation)."""
    import embeddings
    try:
        import numpy as np
    except ImportError:
        return None                                     # no numpy -> caller falls back to TF-IDF (honest degradation)
    q_tf = Counter(_tokenize(query))
    if not q_tf:
        return []                                       # query with no tokens -> no results (same as TF-IDF)
    q_emb = embeddings.embed_query(query)
    if q_emb is None:
        return None
    q = np.asarray(q_emb, dtype="float32")
    idf = index["idf"]
    q_vec = {t: f * idf.get(t, 0.0) for t, f in q_tf.items()}
    q_norm = math.sqrt(sum(v * v for v in q_vec.values())) or 1.0
    filt = (only_source or "").lower()
    scored = []
    for ch in index["chunks"]:
        if filt and filt not in str(ch.get("source", "")).lower():
            continue                                    # restricted to the mentioned file
        lex = 0.0
        dot = sum(q_vec.get(t, 0.0) * (f * idf.get(t, 0.0)) for t, f in ch["tf"].items())
        if dot > 0:
            lex = dot / (q_norm * ch.get("norm", 1.0))
        sem = 0.0
        emb = ch.get("emb")
        if emb is not None:
            try:                                        # a corrupt emb (odd dimension) only cancels ITS
                v = np.asarray(emb, dtype="float32")    # semantic signal; doesn't break the whole search
                if v.shape == q.shape:
                    sem = embeddings.rescale_semantic(float(np.dot(q, v)))
            except Exception:  # noqa: BLE001
                sem = 0.0
        score = embeddings.W_LEX * lex + embeddings.W_SEM * sem
        if score > 0:
            scored.append((score, ch))
    scored.sort(key=lambda x: x[0], reverse=True)
    return [{"source": ch["source"], "line": ch.get("line", 1),
             "text": ch["text"], "score": round(s, 4)} for s, ch in scored[:max(1, k)]]


def search(base, query, k=4, only_source=None, cfg=None):
    """Return the TOP-K most relevant fragments for `query`: list of {source, line, text, score}.
    (Re)indexes if needed. Never returns whole documents -> light context.
    `only_source`: if given, RESTRICTS the search to the fragments of that file (e.g. when the
    question names a specific document: 'according to my CV' -> only the CV).
    If the embeddings layer is ON (and the model available), the search is HYBRID (lexical +
    semantic -> cross-language). With OFF/degraded, it's the usual TF-IDF: if the file doesn't match
    lexically (e.g. question in ES, document in EN), it returns [] -> honest cutoff, no making things up."""
    base = Path(base)
    index_docs(base, cfg=cfg)                           # ensure a fresh index (with embeddings if the layer is ON)
    index = _load_index(base)
    if not index or not index.get("chunks"):
        return []
    if index.get("emb_model"):                          # hybrid only if the index carries embeddings...
        import embeddings
        if embeddings.enabled(cfg) and embeddings.available(cfg):   # ...and the layer is still active and available
            r = _hybrid_search(index, query, only_source, k)
            if r is not None:
                return r                                # None -> embedding the query failed -> fall back to TF-IDF
    # --- Pure TF-IDF (the usual behavior, IDENTICAL with the layer OFF) ---
    idf = index["idf"]
    q_tf = Counter(_tokenize(query))
    if not q_tf:
        return []
    q_vec = {t: f * idf.get(t, 0.0) for t, f in q_tf.items()}
    q_norm = math.sqrt(sum(v * v for v in q_vec.values())) or 1.0
    filt = (only_source or "").lower()
    scored = []
    for ch in index["chunks"]:
        if filt and filt not in str(ch.get("source", "")).lower():
            continue                                    # restricted to the mentioned file
        dot = sum(q_vec.get(t, 0.0) * (f * idf.get(t, 0.0)) for t, f in ch["tf"].items())
        if dot > 0:
            scored.append((dot / (q_norm * ch.get("norm", 1.0)), ch))
    scored.sort(key=lambda x: x[0], reverse=True)
    return [{"source": ch["source"], "line": ch.get("line", 1),
             "text": ch["text"], "score": round(s, 4)} for s, ch in scored[:max(1, k)]]


def format_result(results):
    """Format the top-k for the model/CLI, citing the source (light context: only the fragments)."""
    if not results:
        return "(found nothing relevant in the project's documents)"
    parts = []
    for r in results:
        parts.append(f"[{r['source']}:{r['line']}] {r['text'][:CHUNK_CHARS]}")   # full fragment (was 900<1000)
    return "\n\n".join(parts)
