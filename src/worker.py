"""src/worker.py, Level 15: WORKER MODE (agentic) in a tab separate from the chat.

The worker "works on the user's behalf" inside a PROJECT: it takes a task, breaks it down and runs it,
emitting the steps LIVE. Additive: the chat is untouched; this lives in its own view (`/agentic`) and its
own WebSocket (`/ws-agentic`).

PHASE 1, STUB engine (fake): emits a few example steps to test the UI end to end, WITHOUT touching files
or the model. The REAL engine (reusing the agent) arrives in Phase 2.

Event protocol (consumed by `/ws-agentic` and `web/agentic.html`):
    {type: "worker_start", task, project}
    {type: "worker_plan",  steps: [str, ...]}          # subtasks
    {type: "worker_step",  i: int, state: "pending"|"in_progress"|"done"|"stopped"}
    {type: "worker_log",   text}                       # live activity line
    {type: "worker_final", state: "done"|"stopped"|"error", summary}
"""
import ast
import csv
import html.parser
import json
import os
import re
import subprocess  # nosec B404: runs the TEST of the created script (pytest, fixed list, no shell, with timeout)
import sys
import time
import unicodedata
import zipfile
from pathlib import Path

_STUB_STEPS = [
    "Analyze the task and the project",
    "Prepare a plan of subtasks",
    "Run the subtasks (simulated)",
    "Verify the result (simulated)",
]


def run_stub(task, project, emit, pause=0.4):
    """Fake engine (Phase 1): emits example steps to validate the UI end to end. NEVER touches files or the
    model. `emit(ev: dict)` is the callback (the WebSocket queues it thread-safely). `pause` is the seconds
    between steps (0 in tests, so they run instantly)."""
    emit({"type": "worker_start", "task": task, "project": project})
    emit({"type": "worker_plan", "steps": list(_STUB_STEPS)})
    for i, p in enumerate(_STUB_STEPS):
        emit({"type": "worker_step", "i": i, "state": "in_progress"})
        if pause:
            time.sleep(pause)
        emit({"type": "worker_log", "text": f"(stub) {p} - ok"})
        emit({"type": "worker_step", "i": i, "state": "done"})
    emit({"type": "worker_final", "state": "done",
          "summary": "Fake engine (Phase 1): UI flow tested end to end."})


# Phase 2: REAL worker engine (reuses the agent) + DETERMINISTIC acceptance.

_IGNORE_EXTS = (".snapshots", ".knowledge", ".trash", "chats", "__pycache__")


def verify_static_web(proj_dir, since=None):
    """DETERMINISTIC acceptance of the "single-page static web" MVP task. By CODE, not by the model's
    judgment. Checks the project folder for index.html + one .css; and that the index has a doctype,
    <title>, body content, a LINK to the CSS, and that the HTML PARSES without errors. `since` (mtime):
    only credits files created/modified IN THIS run (not a PREEXISTING index.html/css, so no false pass).
    Returns (ok: bool, details: list[(bool, str)])."""
    d = Path(proj_dir)
    def _from_this_run(p):
        return p.is_file() and (since is None or p.stat().st_mtime >= since)
    index = d / "index.html"
    css = [p for p in d.glob("*.css") if _from_this_run(p)]
    checks = []
    ok_index = _from_this_run(index)
    checks.append((ok_index, "index.html was created in this run"))
    checks.append((len(css) > 0, "a .css stylesheet was created in this run"))
    content = index.read_text(encoding="utf-8", errors="replace") if ok_index else ""
    low = content.lower()
    checks.append(("<!doctype html" in low, "index.html has <!doctype html>"))
    checks.append(("<title" in low and "</title>" in low, "index.html has <title>...</title>"))
    links = bool(re.search(r'<link[^>]+href=["\'][^"\']+\.css["\']', content, re.I))
    checks.append((links, "index.html links the CSS (<link ... .css>)"))
    visible_text = re.sub(r"<[^>]+>", "", content).strip()
    checks.append(("<body" in low and len(visible_text) >= 10, "index.html has content in the body"))
    ok_parse = bool(content)
    try:                                    # parses without blowing up (html.parser is lenient but catches gross garbage)
        p = html.parser.HTMLParser()
        p.feed(content)
        p.close()
    except Exception:                       # noqa: BLE001
        ok_parse = False
    checks.append((ok_parse, "the HTML parses without errors"))
    return all(c for c, _ in checks), checks


def _translator(emit):
    """Returns (tr, state): `tr(ev)` translates AGENT events (plan/step/tool/verdict/hint...) into
    `worker_*` events for the agentic view, keeping track of step state. Other events
    (delta/brains/memory/intent...) are ignored (noise for the worker view)."""
    state = {"n": 0, "current": -1}

    def tr(ev):
        t = ev.get("type")
        if t == "plan":
            steps = ev.get("steps") or []
            state["n"] = len(steps)
            emit({"type": "worker_plan", "steps": steps})
        elif t == "step":
            i = int(ev.get("i", 0)) - 1                 # the chat numbers 1..N; the worker view 0..N-1
            if state["current"] >= 0 and state["current"] < i:
                emit({"type": "worker_step", "i": state["current"], "state": "done"})
            if 0 <= i < state["n"]:
                emit({"type": "worker_step", "i": i, "state": "in_progress"})
                state["current"] = i
        elif t == "tool":
            name = ev.get("name") or "tool"
            emit({"type": "worker_log", "text": f"→ {name} {str(ev.get('args') or '')[:110]}"})
        elif t == "verdict":
            emit({"type": "worker_log", "text": f"verification: {ev.get('state', '')} {str(ev.get('detail') or '')[:110]}"})
        elif t in ("hint", "verifying", "self_critique"):
            emit({"type": "worker_log", "text": str(ev.get("text") or t)[:160]})

    return tr, state


def classify_task(task):
    """DETERMINISTIC classification of the task type (decides the verifier and the guarantee level):
    'python_test' (script + test); Phase B: 'csv', 'ipynb', 'xlsx', 'pptx', 'docx' (each with its own
    verifier that OPENS/parses the file); 'web' (static web); 'open' (attempted anyway, with generic
    checks and an honest status, never a false pass)."""
    low = (task or "").lower()
    is_py = "python" in low or ".py" in low
    is_test = "test" in low or "test" in low or "pytest" in low
    if is_py and is_test:
        return "python_test"
    # 1) Explicit EXTENSION -> unambiguous class (wins over everything else).
    if ".csv" in low:
        return "csv"
    if ".ipynb" in low:
        return "ipynb"
    if ".xlsx" in low:
        return "xlsx"
    if ".pptx" in low:
        return "pptx"
    if ".docx" in low:
        return "docx"
    # 2) WEB deliverable signal -> 'web' BEFORE the office topic words, so "make a web about excel" /
    #    "landing about a spreadsheet" is not stolen toward the file.
    if any(k in low for k in ("web", "page", "page", "html", "site", "landing")):
        return "web"
    # 3) File-type keywords, with WORD BOUNDARY (avoids 'excellent'->xlsx, 'representation'->pptx) and
    #    disambiguating 'notebook' (in Latin America = laptop).
    if re.search(r"\bcsv\b", low):
        return "csv"
    if "jupyter" in low or (re.search(r"\bnotebook\b", low)
                            and not any(k in low for k in ("laptop", "laptop", "laptop", "lenovo",
                                                           "vender", "buy", "shop", "computer"))):
        return "ipynb"
    if re.search(r"\bexcel\b", low) or "spreadsheet" in low or "spreadsheet" in low:
        return "xlsx"
    if ("powerpoint" in low or "power point" in low
            or re.search(r"\bpresentations?\b", low) or re.search(r"\bslides?\b", low)):
        return "pptx"
    if re.search(r"\bword\b", low) or "word document" in low:
        return "docx"
    return "open"


def _verifier_for(cls, cfg, task, since=None):
    """Returns the deterministic verifier (proj_dir -> (ok, details)) for a class, or None if the class is
    'open' (no strong verifier -> honest status). `since` (mtime) requires the Phase B verifiers to confirm
    the file was created IN THIS run (don't credit a preexisting one, so no false pass)."""
    if cls == "web":
        return lambda d: verify_static_web(d, since)
    if cls == "python_test":
        return lambda d: verify_python_test(d, cfg, since)
    if cls == "csv":
        return lambda d: verify_csv(d, task, since)
    if cls == "ipynb":
        return lambda d: verify_ipynb(d, cfg, since)
    if cls in ("docx", "xlsx", "pptx"):
        return lambda d: verify_office(d, cls, since)
    return None


def verify_python_test(proj_dir, cfg=None, since=None):
    """DETERMINISTIC acceptance of "Python script + its test": a .py script exists, its test exists
    (test_*.py / *_test.py) and the test PASSES by running pytest in a BOUNDED way, inside the project
    (cwd), with a time limit (`WORKER_TEST_MAX_SEC`), no shell and no proxy variables. `since` (mtime):
    only counts .py files created/modified IN THIS run (not a PREEXISTING script+test, so no false pass).
    HONEST LIMIT: there is no hard network/process isolation (same level as the agent's run_bash).
    Returns (ok, details) like the other verifiers."""
    d = Path(proj_dir)
    def _from_this_run(p):
        return p.is_file() and (since is None or p.stat().st_mtime >= since)
    tests = sorted(set(p for p in list(d.glob("test_*.py")) + list(d.glob("*_test.py")) if _from_this_run(p)))
    scripts = [p for p in d.glob("*.py") if _from_this_run(p) and p not in tests]
    checks = [(len(scripts) > 0, "a .py script exists (besides the test)"),
              (len(tests) > 0, "its test exists (test_*.py)")]
    if not tests or not scripts:
        checks.append((False, "the test cannot be run (the script or the test is missing)"))
        return False, checks
    limit = int(_cfg_num(cfg or {}, "WORKER_TEST_MAX_SEC", 60))
    env = {k: v for k, v in os.environ.items()
           if k.upper() not in ("HTTP_PROXY", "HTTPS_PROXY", "ALL_PROXY", "FTP_PROXY")}
    # N18 B1: the message reports the real DETERMINISTIC RESULT (pass/fail), NOT the model's claim. Before,
    # the text always said "the test PASSES" (even with returncode != 0) -> "the test PASSES (pytest: 1
    # failed)", a contradiction. Consistent with "deterministic over the model's judgment".
    try:
        r = subprocess.run(  # nosec B603: venv interpreter + pytest, fixed args, confined cwd, timeout
            [sys.executable, "-m", "pytest", "-q", "--maxfail=1", "-p", "no:cacheprovider",
             *[t.name for t in tests]],
            cwd=str(d), capture_output=True, text=True, timeout=limit, env=env)
        tail = ((r.stdout or "").strip().splitlines() or ["(no output)"])[-1][:120]
        passed = r.returncode == 0
        checks.append((passed, f"the test {'PASSES' if passed else 'FAILS'} (pytest: {tail})"))
    except subprocess.TimeoutExpired:
        checks.append((False, f"the test FAILS (pytest did not finish in {limit}s; cancelled)"))
    except Exception as e:  # noqa: BLE001: not being able to run the test = not verified (honest)
        checks.append((False, f"the test FAILS (could not run it: {type(e).__name__})"))
    return all(c for c, _ in checks), checks


# Phase B: the worker CREATES deliverables (csv, ipynb, docx/xlsx/pptx) with a DETERMINISTIC verifier per
# type. Deterministic over the model's judgment: the file is actually opened/parsed; whatever doesn't pass
# -> honest status (review/stopped), NEVER a false pass.
def _most_recent(paths, since=None):
    """The most recently modified file in `paths` (the worker may leave several). If `since` (mtime) is
    given, IGNORE files older than that instant: this way the verifier only credits a file created/modified
    DURING this run, never a preexisting one (avoids a false pass over someone else's work)."""
    ps = [p for p in paths if p.is_file() and (since is None or p.stat().st_mtime >= since)]
    return max(ps, key=lambda p: p.stat().st_mtime) if ps else None


def _strip_accents(s):
    """lowercase + accent-stripped (to compare column names ignoring accents, since models often write
    headers without the accent)."""
    s = (s or "").strip().lower()
    return "".join(c for c in unicodedata.normalize("NFD", s) if unicodedata.category(c) != "Mn")


# Columns requested ONLY given an UNAMBIGUOUS enumeration: the keyword followed by ':' ("columns: a, b,
# c"). Without ':' it is too ambiguous ("the electromagnetic fields") and produced false negatives.
_RX_COLUMNS = re.compile(r"(?:columns?|fields?|headers?)\s*:\s*([^.\n]+)", re.IGNORECASE)
_NOT_A_COLUMN = ("row", "record", "example", "data", "table", "column", "field",
                 "header", "line", "sheet", "cell")


def _requested_columns(task):
    """If the task ENUMERATES columns after a colon ('columns: name, age, city'), extract them to
    require them. VERY conservative: only "clean" names (1-2 words, no digits or instruction words); when
    in doubt discard the fragment. Without ':' -> empty list (structural check only)."""
    m = _RX_COLUMNS.search(task or "")
    if not m:
        return []
    cols = []
    for p in re.split(r"\s*,\s*|\s*;\s*|\s+y\s+|\s*/\s*", m.group(1)):
        c = p.strip(" '\"«».:").strip().lower()
        if (not c or len(c) > 25 or ":" in c or re.search(r"\d", c)
                or len(c.split()) > 2 or any(s in c for s in _NOT_A_COLUMN)):
            continue                                    # it's an instruction or noise, not a column
        cols.append(c)
    return cols[:12]


def verify_csv(proj_dir, task="", since=None):
    """DETERMINISTIC acceptance of a .csv created in this run: it exists, PARSES with the csv module, has a
    header of >=2 columns + >=1 data row with a CONSISTENT column count, and, if the task enumerates them
    after ':', contains the requested COLUMNS (compared accent-insensitively). >=2 columns avoids accepting
    comma-less prose as a "1-column csv"."""
    d = Path(proj_dir)
    f = _most_recent(d.glob("*.csv"), since)
    checks = [(f is not None, "a .csv was created in this run")]
    if f is None:
        checks.append((False, "the .csv parses with header and data"))
        return False, checks
    try:
        text = f.read_text(encoding="utf-8-sig", errors="replace")     # utf-8-sig: drops the BOM if present
        rows = [r for r in csv.reader(text.splitlines()) if any((c or "").strip() for c in r)]
    except Exception as e:  # noqa: BLE001
        checks.append((False, f"the .csv parses ({type(e).__name__})"))
        return False, checks
    ncol = len(rows[0]) if rows else 0
    ok_struct = len(rows) >= 2 and ncol >= 2 and all(len(r) == ncol for r in rows)
    checks.append((ok_struct, f"{f.name} parses: {len(rows)} rows x {ncol} columns (>=2), consistent"))
    requested = _requested_columns(task)
    if requested:
        header = {_strip_accents(c) for c in (rows[0] if rows else [])}
        missing = [c for c in requested if _strip_accents(c) not in header]
        checks.append((not missing, "it has the requested columns"
                       + (f" (missing: {', '.join(missing)})" if missing else f" ({', '.join(requested)})")))
    return all(c for c, _ in checks), checks


def _has_real_code(src):
    """True if `src` has >=1 truly EXECUTABLE statement. Comments vanish in the AST, so a '# TODO:
    implement' gives an empty body -> False; we also discard 'pass' and a lone docstring/literal (they
    don't "do" anything). Avoids accepting a stub notebook. If it doesn't parse, the separate syntax check
    catches it -> return False here."""
    try:
        body = ast.parse(src or "").body
    except SyntaxError:
        return False
    for node in body:
        if isinstance(node, ast.Pass):
            continue
        if isinstance(node, ast.Expr) and isinstance(getattr(node, "value", None), ast.Constant):
            continue                       # docstring or lone literal: not "doing something"
        return True
    return False


def verify_ipynb(proj_dir, cfg=None, since=None):
    """DETERMINISTIC acceptance of a .ipynb created in this run: it is a VALID NOTEBOOK JSON (dict with
    'cells' and 'nbformat'), has >=1 code cell WITH REAL CODE (not empty, not only magics) and NO code cell
    with a SYNTAX error (compile, running nothing). Robust to malformed cells/source (never propagates an
    exception: that would take down the pass)."""
    d = Path(proj_dir)
    f = _most_recent(d.glob("*.ipynb"), since)
    checks = [(f is not None, "an .ipynb notebook was created in this run")]
    if f is None:
        checks.append((False, "the notebook is valid JSON and its cells compile"))
        return False, checks
    try:
        data = json.loads(f.read_text(encoding="utf-8", errors="replace"))
        struct_ok = isinstance(data, dict) and isinstance(data.get("cells"), list) and "nbformat" in data
    except Exception:  # noqa: BLE001
        struct_ok, data = False, None
    checks.append((struct_ok, f"{f.name} is a valid notebook JSON (cells + nbformat)"))
    if not struct_ok:
        return False, checks
    code = [c for c in data["cells"] if isinstance(c, dict) and c.get("cell_type") == "code"]
    n_with_code, n_bad, first_err = 0, 0, ""
    for c in code:
        src = c.get("source")
        if isinstance(src, list):
            src = "".join(s for s in src if isinstance(s, str))   # cells with malformed source don't crash
        elif not isinstance(src, str):
            src = ""
        src = re.sub(r"^\s*[%!].*$", "", src, flags=re.MULTILINE)   # ignore Jupyter magics/shell (%..., !...)
        if _has_real_code(src):        # >=1 real executable statement (a '# TODO' or 'pass' does NOT count)
            n_with_code += 1
        try:
            compile(src, "<cell>", "exec")
        except SyntaxError as e:
            n_bad += 1
            first_err = first_err or f"{e.msg} (line {e.lineno})"
    checks.append((n_with_code > 0, "it has at least one cell with real code (not empty, not only comments)"))
    checks.append((n_bad == 0 and n_with_code > 0, "the code cells have no syntax errors"
                   + (f" - failed: {first_err}" if n_bad else "")))
    return all(c for c, _ in checks), checks


_OOXML_PART = {"docx": "word/document.xml", "xlsx": "xl/workbook.xml", "pptx": "ppt/presentation.xml"}


def _ooxml_has_content(z, kind):
    """Does the OOXML have REAL content, not just empty boilerplate? Without a library, reading the zip:
    docx -> a non-empty <w:t> text run; xlsx -> some <c> cell in some sheet; pptx -> >=1 slide with real
    TEXT <a:t> (an empty slide isn't enough). (XML size isn't enough: an empty document already exceeds the
    boilerplate byte count.)"""
    names = z.namelist()
    try:
        if kind == "docx":
            return bool(re.search(rb"<w:t[ >][^<]*\S", z.read("word/document.xml")))
        if kind == "xlsx":
            sheets = [n for n in names if re.match(r"xl/worksheets/sheet\d+\.xml$", n)]
            return any(re.search(rb"<c[ >]", z.read(n)) for n in sheets)
        if kind == "pptx":
            slides = [n for n in names if re.match(r"ppt/slides/slide\d+\.xml$", n)]
            return any(re.search(rb"<a:t[ >][^<]*\S", z.read(n)) for n in slides)   # slide WITH text
    except Exception:  # noqa: BLE001
        return False
    return False


def _validate_ooxml_zip(f, kind):
    """Library-free fallback: a docx/xlsx/pptx is an OOXML ZIP, so we validate that the zip is not corrupt,
    contains its main part and has REAL content (not just empty boilerplate). Detects a corrupt / empty /
    renamed file."""
    try:
        with zipfile.ZipFile(str(f)) as z:
            if z.testzip() is not None:
                return False, "corrupt zip (CRC)"
            if _OOXML_PART[kind] not in z.namelist():
                return False, f"missing {_OOXML_PART[kind]} (not a valid {kind})"
            if not _ooxml_has_content(z, kind):
                return False, "OOXML with no real content (empty document)"
            return True, "valid OOXML with content"
    except Exception as e:  # noqa: BLE001
        return False, f"not a valid OOXML: {type(e).__name__}"


def _open_office(f, kind):
    """Opens the file with ITS library (openpyxl/python-pptx/python-docx) if present, a strong check that
    it isn't corrupt and has content; if the library is missing, falls back to the zip OOXML validator."""
    try:
        if kind == "xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(str(f), read_only=True, data_only=True)
            try:
                sheets = wb.sheetnames
                rows = sum(1 for _ in wb[sheets[0]].iter_rows()) if sheets else 0
            finally:
                wb.close()
            return (len(sheets) >= 1 and rows >= 1, f"{len(sheets)} sheet(s), {rows} row(s)")
        if kind == "pptx":
            from pptx import Presentation
            pres = Presentation(str(f))
            n = len(pres.slides)
            with_text = sum(1 for s in pres.slides for sh in s.shapes
                            if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip())
            return (n >= 1 and with_text >= 1, f"{n} slide(s), {with_text} with text")
        if kind == "docx":
            import docx  # python-docx (if present)
            doc = docx.Document(str(f))
            text = "".join(p.text for p in doc.paragraphs).strip()
            return (len(text) >= 1, f"{len(doc.paragraphs)} paragraph(s)")
    except ImportError:
        return _validate_ooxml_zip(f, kind)          # no library -> zip validation (honest)
    except Exception as e:  # noqa: BLE001: doesn't open = corrupt = not verified
        return False, f"doesn't open ({type(e).__name__})"
    return _validate_ooxml_zip(f, kind)


def verify_office(proj_dir, kind, since=None):
    """DETERMINISTIC acceptance of a docx/xlsx/pptx created in this run: it exists and OPENS WITHOUT
    CORRUPTION with content (using its library; or, if missing, validating the zip's OOXML structure).
    kind in {docx,xlsx,pptx}."""
    d = Path(proj_dir)
    ext = "." + kind
    f = _most_recent(d.glob(f"*{ext}"), since)
    checks = [(f is not None, f"a {ext} was created in this run")]
    if f is None:
        checks.append((False, f"the {ext} opens without corruption"))
        return False, checks
    ok, detail = _open_office(f, kind)
    checks.append((ok, f"{f.name} opens without corruption and has content ({detail})"))
    return all(c for c, _ in checks), checks


def _useful_files(proj_dir):
    """Lightweight snapshot of the project's "working" files (excludes hidden, chats/, __pycache__)."""
    out = set()
    d = Path(proj_dir)
    if not d.is_dir():
        return out
    for p in d.rglob("*"):
        rel = p.relative_to(d)
        if any(part.startswith(".") or part in ("chats", "__pycache__") for part in rel.parts):
            continue
        if p.is_file():
            out.add(str(rel))
    return out


def generic_checks(proj_dir, before):
    """Phase A: GENERIC checks for open tasks (no strong verifier): did it leave new/modified files, and if
    there are new .py files do they compile without error? These do NOT guarantee the task (which is why the
    final status is "done WITHOUT strong verification"). Returns (something_useful, details, new_files)."""
    import py_compile
    d = Path(proj_dir)
    new_files = sorted(_useful_files(d) - set(before or ()))
    details = [(len(new_files) > 0, f"left files in the project ({len(new_files)}: {', '.join(new_files[:4])})"
                if new_files else "left files in the project")]
    for rel in new_files:
        if rel.endswith(".py"):
            try:
                py_compile.compile(str(d / rel), doraise=True)
                details.append((True, f"{rel} compiles without errors"))
            except Exception:  # noqa: BLE001
                details.append((False, f"{rel} compiles without errors"))
    return all(c for c, _ in details) and bool(new_files), details, new_files


# N18 Phase A: CONTRACT-FIRST worker. In a "module + its test" task, every model call is fresh and the 14B
# doesn't reliably remember the interface of the code it wrote earlier (names, signatures, return type), so
# the test doesn't match the module and pytest fails even when the logic is correct. Fix: RE-READ before
# depending. (1) the initial spec asks to create the module, RE-READ it, then the exact test; (2)
# deterministic and strong: on each RETRY the worker re-reads the REAL module from disk and injects it into
# the context, so the model doesn't have to remember the interface: it SEES it. Config-driven and
# reversible (WORKER_CONTRACT_FIRST); no changes to the chat.
def _contract_first(cfg):
    """Is "contract-first" mode enabled? (WORKER_CONTRACT_FIRST). Adopted ON in the template
    (config.env.example) after the N18 hard bench (OFF 50% vs ON 67%). The code FALLBACK, if the key is
    missing, is OFF (conservative: the evidence is moderate, not decisive); config.env turns it ON.
    Reversible by config."""
    return str((cfg or {}).get("WORKER_CONTRACT_FIRST", "0")).strip().lower() not in ("0", "false", "no", "off", "")


_CONTRACT_MAX_CHARS = 1800   # per file (light context: the interface fits easily)


def _contract_test_block(proj_dir, since=None):
    """Re-reads the real MODULES (.py that are NOT tests) created in this run and returns a text block with
    their content, to inject when writing/fixing the test -> the model matches the test against the REAL
    interface, not against what it thinks it remembers. '' if there is no module. Deterministic (reads
    disk, no opinion)."""
    d = Path(proj_dir)
    def _is_test(p):
        return p.name.startswith("test_") or p.name.endswith("_test.py")
    modules = sorted(p for p in d.glob("*.py")
                     if p.is_file() and not _is_test(p) and (since is None or p.stat().st_mtime >= since))
    if not modules:
        return ""
    parts = []
    for p in modules[:2]:                                    # at most 2 modules (light context)
        try:
            txt = p.read_text(encoding="utf-8", errors="replace")[:_CONTRACT_MAX_CHARS]
        except OSError:
            continue
        parts.append(f"# ---- {p.name} (REAL module, already on disk) ----\n{txt}")
    if not parts:
        return ""
    return ("\n\nREAL CONTRACT: this is the module(s) you ALREADY wrote. Write/fix the test so it MATCHES "
            "this interface EXACTLY: same import names, same function signatures and the same return type. "
            "Do NOT assume the interface: use it exactly as it is here.\n\n"
            + "\n\n".join(parts))


def _task_for_agent(task, cls=None, contract_first=False):
    """Worker DECOMPOSITION: elaborates the task into a CONCRETE spec by class, so a small model can fulfill
    it in one go. It is not "cheating" the verifier: it translates the request into what the complete result
    IS. Open tasks pass THROUGH unchanged."""
    if cls is None:
        cls = classify_task(task)
    if cls == "web":
        return (task.strip() + "\n\nBe specific (everything in the current folder): create TWO files: "
                "(1) `index.html`: a complete HTML page with `<!doctype html>`, a `<title>`, visible "
                "content inside the `<body>`, and that LINKS the stylesheet with "
                "`<link rel=\"stylesheet\" href=\"styles.css\">` in the `<head>`; "
                "(2) `styles.css`: some basic styles. Make sure the `index.html` keeps the `<link>` to "
                "`styles.css`. Don't open the browser or run commands: just create the files.")
    if cls == "python_test":
        contract = ("" if not contract_first else
                    " IMPORTANT (CONTRACT-FIRST order): create the module script FIRST; THEN RE-READ it "
                    "with `read_file` to see its REAL interface (imports, function names and signatures, "
                    "return type); and ONLY AFTER THAT write the test, matching it EXACTLY to that real "
                    "interface (not the one you think you remember).")
        return (task.strip() + "\n\nBe specific (everything in the current folder): create TWO files: "
                "(1) a Python script with the requested functionality; (2) its test `test_<name>.py` "
                "(pytest style: `test_*` functions with asserts, importing the script) covering the "
                "important cases." + contract + " Run `python -m pytest -q` and if any test fails, fix it "
                "until they ALL pass. Don't use the network or open anything: just the files and pytest.")
    # Phase B: file deliverables. Reuses the chat tools; be CONCRETE.
    if cls == "csv":
        return (task.strip() + "\n\nBe specific (in the current folder): create ONE `.csv` file with "
                "`write_file`. First line = header with the column names; one row per record, values "
                "separated by commas and the SAME number of columns in every row. Include at least one "
                "data row. Don't open anything or use the network.")
    if cls == "ipynb":
        return (task.strip() + "\n\nBe specific (in the current folder): create ONE `.ipynb` notebook with "
                "the `create_notebook` tool (cells = list of {type, source}; type 'text' for markdown, "
                "'code' for Python). Include AT LEAST one code cell, and make ALL the code valid (no "
                "syntax errors). Don't run the notebook or use the network.")
    if cls == "docx":
        return (task.strip() + "\n\nBe specific (in the current folder): create ONE `.docx` document with "
                "the `create_document` tool (pass it the content in Markdown and the path with the `.docx` "
                "extension; it is converted with pandoc). Give it real content. Don't open anything or use "
                "the network.")
    if cls == "xlsx":
        return (task.strip() + "\n\nBe specific (in the current folder): create the `.xlsx` with the "
                "`create_office_doc` tool, passing it `content={\"rows\": [[<header>], [<data>], ...]}` "
                "(the first row is the header). Do NOT write or run any script. Don't use the network.")
    if cls == "pptx":
        return (task.strip() + "\n\nBe specific (in the current folder): create the `.pptx` with the "
                "`create_office_doc` tool, passing it `content={\"slides\": [{\"title\": \"...\", "
                "\"puntos\": [\"...\", \"...\"]}, ...]}` (at least one slide with title and text). Do NOT "
                "write or run any script. Don't use the network.")
    return task


def _single_pass(task, project, emit, cfg, router, resolve, verifier, proj_dir, contract_first=False):
    """ONE agent pass over the task + deterministic acceptance. Returns (ok, details, error). Reused by
    run_real (Phase 2, one pass) and run_autonomous (Phase 3, loop of passes). `contract_first` (N18):
    guides the "module + test" spec to re-read the real interface (additive, default off)."""
    from events import use_emitter
    tr, state = _translator(emit)
    error = None
    try:
        with use_emitter(tr):
            resolve(router, cfg, project,
                    _task_for_agent(task, contract_first=contract_first))   # the agent creates the files (WORKDIR)
    except Exception as e:                              # noqa: BLE001: never propagate: reported honestly
        error = f"{type(e).__name__}: {e}"
    if state["current"] >= 0:
        emit({"type": "worker_step", "i": state["current"], "state": "done"})
    ok, details = verifier(proj_dir)                # DETERMINISTIC acceptance (by code)
    for c, txt in details:
        emit({"type": "worker_log", "text": f"{'✓' if c else '✗'} {txt}"})
    return ok, details, error


def run_real(task, project, emit, cfg, router, resolve=None, verifier=None):
    """REAL SINGLE-PASS engine (Phase 2): reuses the AGENT + deterministic acceptance. STOPS HONESTLY: only
    "done" if the verifier passes; otherwise "stopped", saying what's missing. `resolve`/`verifier`
    injectable (tests). The lock/WORKDIR are managed by the caller (web_server)."""
    if resolve is None:
        resolve = _default_resolver            # aware of manual folders (project_dir)
    if verifier is None:
        verifier = verify_static_web
    emit({"type": "worker_start", "task": task, "project": project})
    proj_dir = project_dir(cfg, project)     # REAL folder (exact if it exists; slug otherwise)
    ok, details, error = _single_pass(task, project, emit, cfg, router, resolve, verifier, proj_dir)
    if error and not ok:
        emit({"type": "worker_final", "state": "error", "summary": f"The agent failed: {error}"})
    elif ok:
        emit({"type": "worker_final", "state": "done",
              "summary": "Single-page static web created and verified (deterministic acceptance)."})
    else:
        missing = [txt for c, txt in details if not c]
        emit({"type": "worker_final", "state": "stopped",
              "summary": "Not calling it done (honest): missing " + "; ".join(missing[:3]) + "."})


# Phase 3: robust autonomy, a loop of verified passes + limits + anti-loop + pause/resume/cancel +
# persistent state + run log.

def project_dir(cfg, name):
    """REAL project folder for agentic mode. If the projects base folder has a subfolder with that EXACT
    name (including ones placed BY HAND by the user: spaces, uppercase, accents...), use THAT; otherwise the
    usual slug path (compatible with the app's projects). PATH-SAFETY: name only (no separators/`..`),
    never hidden, and CONFINED to the base (a symlink pointing outside is ignored and falls back to the
    slug path)."""
    import projects
    base = projects.base_dir(cfg)
    clean = Path(str(name or "")).name.strip()          # name only: kills '../', '/x', etc.
    if clean and not clean.startswith("."):
        d = base / clean
        try:
            if d.is_dir() and d.resolve().parent == base.resolve():   # confined (anti symlink escape)
                return d
        except OSError:
            pass
    return projects.project_path(cfg, clean)


def _default_resolver(router, cfg, project, task):
    """The worker's real resolver. Names that SURVIVE the slug (the app's projects) -> the usual path
    (`resolve_in_project`: per-project memory/chats/snapshot). Folders placed BY HAND (spaces, accents...
    that the slug would alter) -> the agent works INSIDE that EXACT folder: the WORKDIR is set to it (the
    caller already holds the lock and restores the WORKDIR) and the normal pipeline runs. With the usual
    safety: tools stay confined to the WORKDIR (path-safety), trash and limits intact."""
    import projects
    import orchestrator
    d = project_dir(cfg, project)
    # Phase A: the worker ALWAYS goes STRAIGHT to the pipeline (plan->run->verify). `direct=True` skips the
    # chat shortcuts (docs -> skills -> CHAT triage), which reply with TEXT and no tools and were hijacking
    # worker tasks ("write X to a file" -> chat; long specs -> skill).
    if d == projects.project_path(cfg, project):        # "normal" name -> usual behavior
        return orchestrator.resolve_in_project(router, cfg, project, task, direct=True)
    orchestrator._set_workdir(d)                          # manual folder: work INSIDE it
    return orchestrator.resolve_task(router, cfg, task, direct=True)


def _worker_dir(cfg, project):
    """The worker's state folder inside the project (.worker/). Created if needed."""
    d = project_dir(cfg, project) / ".worker"
    d.mkdir(parents=True, exist_ok=True)
    return d


# Worker signals. Phase 3: pause/resume/cancel (between passes). Phase C: approve/redo (between STAGES,
# when the worker is waiting for the user's approval).
_SIGNALS = ("continue", "pause", "cancel", "approve", "redo")


def read_signal(cfg, project):
    """Control signal set by the view (one of _SIGNALS; 'continue' by default). Never raises."""
    try:
        import json
        f = _worker_dir(cfg, project) / "control.json"
        if f.exists():
            return (json.loads(f.read_text(encoding="utf-8")).get("signal") or "continue")
    except Exception:  # noqa: BLE001
        pass
    return "continue"


def _write_json_atomic(path, data):
    """N17 Phase B: truly ATOMIC write (mkstemp in the same dir + os.replace). Before, these states were
    written with write_text (truncate -> write), and a mid-write crash or a concurrent reader (the signal
    polling, GET /worker-status, the stage RESUME from stages.json) could see empty or half-written JSON;
    with a corrupt stages.json, resume-after-a-crash died silently. May raise (the caller decides whether
    to swallow it)."""
    import json
    import tempfile
    path = Path(path)
    fd, tmp = tempfile.mkstemp(dir=str(path.parent), prefix=path.name + ".", suffix=".tmp")
    try:
        with os.fdopen(fd, "w", encoding="utf-8") as fh:
            fh.write(json.dumps(data, ensure_ascii=False))
        os.replace(tmp, str(path))                       # atomic on the same filesystem
    except Exception:
        try:
            os.unlink(tmp)                               # don't leave orphan temp files if something fails
        except OSError:
            pass
        raise


def set_signal(cfg, project, signal):
    """The view sets the signal (pause/resume->continue/cancel/approve/redo). Atomic write. Never raises."""
    signal = signal if signal in _SIGNALS else "continue"
    try:
        _write_json_atomic(_worker_dir(cfg, project) / "control.json", {"signal": signal})
        return True
    except Exception:  # noqa: BLE001
        return False


def save_state(cfg, project, **kv):
    """Persistent per-project worker state (.worker/state.json), for resuming/reading in the morning.
    Atomic write (N17-B). Never raises."""
    try:
        _write_json_atomic(_worker_dir(cfg, project) / "state.json", kv)
    except Exception:  # noqa: BLE001
        pass


def read_state(cfg, project):
    import json
    try:
        f = _worker_dir(cfg, project) / "state.json"
        if f.exists():
            return json.loads(f.read_text(encoding="utf-8"))
    except Exception:  # noqa: BLE001
        pass
    return {}


_RUN_LOG_MAX_LINES = 400   # honest cap: keeps the latest runs, doesn't grow forever


def run_log(cfg, project, line):
    """Run log (readable by the user): .worker/log.md, a truly BOUNDED append (N17-B: the comment said
    "bounded" but it grew without a cap; now, past the maximum, the latest _RUN_LOG_MAX_LINES are kept).
    Never raises."""
    try:
        f = _worker_dir(cfg, project) / "log.md"
        with f.open("a", encoding="utf-8") as fh:
            fh.write(f"- {line}\n")
        lines = f.read_text(encoding="utf-8").splitlines(keepends=True)
        if len(lines) > _RUN_LOG_MAX_LINES:           # prune: keep the most recent
            f.write_text("".join(lines[-_RUN_LOG_MAX_LINES:]), encoding="utf-8")
    except Exception:  # noqa: BLE001
        pass


def _wait_if_paused(signal_fn, emit, poll=0.4, pause_limit=1800):
    """If the signal is 'pause', BLOCK until it becomes 'continue' or 'cancel' (or until pause_limit).
    Returns the final signal ('continue'|'cancel'). Control is exercised BETWEEN passes (one agent pass
    can't be cleanly cancelled mid-flight; this is documented)."""
    s = signal_fn()
    if s == "pause":
        emit({"type": "worker_log", "text": "⏸️ paused (waiting to resume or cancel)..."})
        t0 = time.time()
        while s == "pause" and time.time() - t0 < pause_limit:
            time.sleep(poll)
            s = signal_fn()
        emit({"type": "worker_log", "text": "▶️ resumed" if s != "cancel" else "⏹️ cancelling..."})
    return "cancel" if s == "cancel" else "continue"


# Phase C: staged execution with the user's approval (human-in-the-loop). The worker does ONE stage, STOPS
# and waits for "approve / redo / cancel". The user is the checkpoint between stages -> it sidesteps the
# local model's ceiling. Persistent state (.worker/stages.json) to RESUME after a crash. Reuses the signals
# and the pause/resume already built.
def _wait_for_approval(cfg, project, signal_fn, emit, poll=0.5, limit=None):
    """BLOCKS after a stage until the user decides: returns 'approve' | 'redo' | 'cancel'. CONSUMES
    the signal (leaves it at 'continue') so the next stage doesn't inherit it. No decision within `limit`
    seconds -> honest 'cancel' (doesn't hang forever)."""
    limit = limit or float(_cfg_num(cfg, "WORKER_APPROVAL_MAX_SEC", 3600))
    # DISCARD any PREVIOUS signal (a stale/early approval must not approve this stage before the user has
    # seen its result): only a decision sent AFTER this notice counts. It's safe: it runs on the worker
    # thread before the view receives the event and can press a button.
    set_signal(cfg, project, "continue")
    emit({"type": "worker_log",
          "text": "⏸️ stage finished - waiting for your approval (Approve / Redo / Cancel)..."})
    t0 = time.time()
    while time.time() - t0 < limit:
        s = signal_fn()
        if s in ("approve", "redo", "cancel"):
            set_signal(cfg, project, "continue")           # consume: the next stage won't inherit it
            return s
        time.sleep(poll)
    emit({"type": "worker_log", "text": "⏹️ no timely response: cancelling this staged task."})
    return "cancel"


def read_stages(cfg, project):
    """Persisted per-STAGE state (.worker/stages.json): {task, stages:[txt], estados:[st], actual}."""
    import json
    try:
        f = _worker_dir(cfg, project) / "stages.json"
        if f.exists():
            return json.loads(f.read_text(encoding="utf-8"))
    except Exception:  # noqa: BLE001
        pass
    return {}


def save_stages(cfg, project, task, stages, states, current, summary=""):
    """Persists the staged progress (to resume after a crash). Atomic write (N17-B: for real; before, a
    mid-write crash left stages.json corrupt and RESUME died silently). Never raises."""
    try:
        data = {"task": task, "stages": list(stages), "estados": list(states),
                "actual": int(current), "summary": (summary or "")[:300]}
        _write_json_atomic(_worker_dir(cfg, project) / "stages.json", data)
    except Exception:  # noqa: BLE001
        pass


# STAGE markers the user writes naturally: arrows, numbered lists, or connectors.
_RX_ARROW = re.compile(r"\s*(?:→|->|=>)\s*")
# "1) ...", "2. ...", "step 3:" in any position (even on a single line); >=2 are required to split.
_RX_NUMBERED = re.compile(r"(?:^|\s)(?:\d+\s*[\).]|step\s+\d+\s*[:\.\-]?)\s+", re.IGNORECASE)
# The connector must be ADVERBIAL (it separates clauses: "..., then do X"), not part of a word.
# Only "then", "next" and "afterwards" split; a leading "and" is optional.
_RX_CONNECTOR = re.compile(
    r"[;.]?\s+(?:and\s+)?(?:then|next|afterwards?)\b[,:]?\s+",
    re.IGNORECASE)


def _split_into_stages(task):
    """DECOMPOSES the task into stages DETERMINISTICALLY, by the markers the user writes: arrows (→ / -> /
    =>), a numbered list (1) ... 2) ...), or ADVERBIAL connectors ("then", "after", "next",
    not the "after DE ..." locution). Returns the list of stages (>=2) or [task] if there is no clear
    decomposition (-> a single stage, normal behavior). The per-batch stage-count cap is applied, HONESTLY
    and with a warning, in run_by_stages, not here (nothing is discarded here)."""
    t = (task or "").strip()
    if not t:
        return [t]
    if _RX_ARROW.search(t):
        parts = _RX_ARROW.split(t)
    elif len(_RX_NUMBERED.findall(t)) >= 2:                # at least two "N)"/"step N" markers
        parts = _RX_NUMBERED.split(t)
    elif _RX_CONNECTOR.search(t):
        parts = _RX_CONNECTOR.split(t)
    else:
        return [t]
    stages = [p.strip(" ,.;:-\n\t") for p in parts if p and p.strip(" ,.;:-\n\t")]
    return stages if len(stages) >= 2 else [t]


def run_autonomous(task, project, emit, cfg, router, resolve=None, verifier=None, signal_fn=None):
    """PHASE 3: loop of VERIFIED passes with: config limits (`WORKER_MAX_STEPS`, `WORKER_MAX_SEC`),
    ANTI-LOOP by a signature of missing checks (progress is measured by DETERMINISTIC acceptance, not by
    reads, so re-reading doesn't count as progress), and pause/resume/cancel (via `signal_fn`, persistent
    per project). STOPS HONESTLY in every case. State + run log persisted per project."""
    if resolve is None:
        resolve = _default_resolver            # aware of manual folders (project_dir)
    # Phase A: the worker is GENERAL: it classifies the task and picks a verifier. With no strong verifier
    # (class "open") it does NOT invent one: it's attempted anyway and ends in a different HONEST status.
    cls = classify_task(task)
    # Run base instant: the Phase B verifiers only credit files created/modified from here on (2s margin
    # for FS mtime resolution) -> never a false pass over a preexisting deliverable (from another task/chat
    # or placed by hand). Captured BEFORE touching the agent.
    base_mtime = time.time() - 2
    if verifier is None:
        verifier = _verifier_for(cls, cfg, task, since=base_mtime)   # None if 'open' (no strong verifier)
    if signal_fn is None:
        def signal_fn():
            return read_signal(cfg, project)

    # "undefined" bug FIX, defense in depth (the WS already validates): with no valid project, reject
    # honestly BEFORE creating any state (.worker/), never phantom folders.
    project = (project or "").strip()
    if not project or project == "undefined":
        emit({"type": "worker_final", "state": "error",
              "summary": "Choose or create a project first (I didn't get a valid project)."})
        return

    # RESTORATION (N15): besides emitting, the last plan and its step states are CAPTURED to persist them at
    # the end (state.json), so the view can restore the last run after a reload.
    capture = {"plan": [], "steps": {}}
    emit_orig = emit

    def emit(ev):  # noqa: A001: tee: capture + re-emit (same name on purpose: everything below uses it)
        if ev.get("type") == "worker_plan":
            capture["plan"] = list(ev.get("steps") or [])[:30]     # bounded (light context)
            capture["steps"] = {}
        elif ev.get("type") == "worker_step":
            capture["steps"][int(ev.get("i", 0))] = ev.get("state") or "pending"
        emit_orig(ev)

    # (the WORKER_MAX_STEPS/SEG limits are recomputed by _run_verified, which runs the loop)
    emit({"type": "worker_start", "task": task, "project": project})
    set_signal(cfg, project, "continue")                    # clear old signals on start
    save_state(cfg, project, state="running", task=task, attempts=0)
    proj_dir = project_dir(cfg, project)     # REAL folder (exact if it exists; slug otherwise)

    t0 = time.time()
    final_state, summary, attempt = _run_verified(
        task, project, emit, cfg, router, resolve, signal_fn, cls, verifier, proj_dir,
        since=base_mtime, contract_since=base_mtime)   # N19: single task -> run and stage coincide
    emit({"type": "worker_final", "state": final_state, "summary": summary})
    save_state(cfg, project, state=final_state, summary=summary, attempts=attempt,
               seg=round(time.time() - t0, 1), task=task[:300],
               plan=capture["plan"],
               steps=[capture["steps"].get(i, "pending") for i in range(len(capture["plan"]))])
    run_log(cfg, project, f"task '{task[:60]}' -> {final_state} · {summary} (passes={attempt}, {round(time.time()-t0)}s)")


_OK_SUMMARY = {"web": "Single-page static web created and VERIFIED (deterministic acceptance).",
               "python_test": "Python script created and its TEST PASSES (deterministic acceptance).",
               "csv": "CSV created and VERIFIED: parses with header and data (deterministic acceptance).",
               "ipynb": "Notebook .ipynb created and VERIFIED: valid JSON and cells with no syntax errors.",
               "docx": "Document .docx created and VERIFIED: opens without corruption, with content.",
               "xlsx": "Excel .xlsx created and VERIFIED: opens without corruption, with content.",
               "pptx": "Presentation .pptx created and VERIFIED: opens without corruption, with content."}


def _run_verified(task, project, emit, cfg, router, resolve, signal_fn, cls, verifier,
                  proj_dir, since=None, contract_since=None):
    """Runs ONE task (or ONE stage) to its TERMINAL state with the loop of verified passes + anti-loop +
    pause/resume/cancel. Does NOT emit worker_final or persist (the caller decides, so stages can be
    chained). Returns (final_state, summary, attempt). HONEST states: never a false pass.
    - `since` (mtime): the VERIFIER only credits what was created from here on (per-stage -> doesn't credit
      an artifact from a previous stage).
    - `contract_since` (N19): the contract-first INJECTION re-reads modules created from here on. The start
      of the RUN (not the stage) is passed -> in the test stage it DOES find the module created in a
      previous stage (fixes the N18 gap between stages). If None, uses `since` (single task)."""
    max_passes = max(1, int(_cfg_num(cfg, "WORKER_MAX_STEPS", 4)))
    max_seconds = float(_cfg_num(cfg, "WORKER_MAX_SEC", 900))
    contract_first = _contract_first(cfg) and cls == "python_test"   # N18: only for "module + test"
    if contract_since is None:
        contract_since = since
    t0 = time.time()
    signatures = []
    final_state, summary, attempt = "stopped", "", 0
    current_task = task
    if verifier is None:
        # OPEN TASK/STAGE (no strong verifier): attempted anyway (one pass) + generic checks. HONEST status:
        # "done_unverified" (review it) or "stopped/error", NEVER a pass.
        attempt = 1
        emit({"type": "worker_log",
              "text": "no strong verifier: I'll try it and leave it for you to review"})
        before = _useful_files(proj_dir)
        _ok, _det, error = _single_pass(task, project, emit, cfg, router, resolve,
                                        lambda d: (True, []), proj_dir)   # no strong checks in the pass
        something_useful, details, new_files = generic_checks(proj_dir, before)
        for c, txt in details:
            emit({"type": "worker_log", "text": f"{'✓' if c else '✗'} {txt}"})
        if error and not new_files:
            final_state, summary = "error", f"Can't do it: the agent failed ({error})."
        elif something_useful:
            final_state = "done_unverified"
            summary = ("Done, WITHOUT strong verification - review it yourself. It left: "
                       + ", ".join(new_files[:5]) + ("..." if len(new_files) > 5 else "") + ".")
        elif new_files:
            final_state, summary = "stopped", ("Not calling it good: it left files but with problems ("
                                              + "; ".join(t for c, t in details if not c)[:120] + ").")
        else:
            final_state, summary = "stopped", "Can't do it: the task left no file in the project."
    else:
        for attempt in range(1, max_passes + 1):
            if _wait_if_paused(signal_fn, emit) == "cancel":
                final_state, summary = "stopped", "Cancelled by the user."
                break
            if time.time() - t0 > max_seconds:
                final_state, summary = "stopped", f"Time limit reached ({max_seconds:.0f}s)."
                break
            emit({"type": "worker_log", "text": f"- pass {attempt}/{max_passes} -"})
            # N19: DETERMINISTIC contract-first: re-read the REAL MODULE from disk and inject it BEFORE
            # EACH pass (including the 1st), from CODE (not relying on the 14B planner to add a "read"
            # step). On the 1st pass of a single task the module doesn't exist yet -> empty block (no
            # change); in the test stage (`contract_since`=start of the run) or on a retry the module is
            # ALREADY on disk -> its real interface is injected and the test matches it.
            pass_task = current_task
            if contract_first:
                block = _contract_test_block(proj_dir, since=contract_since)
                if block:
                    pass_task = current_task + block
                    emit({"type": "worker_log", "text": "🔎 contract-first: re-reading the real module and using it for the test"})
            ok, details, _err = _single_pass(pass_task, project, emit, cfg, router, resolve,
                                             verifier, proj_dir, contract_first=contract_first)
            if ok:
                final_state = "done"
                summary = _OK_SUMMARY.get(cls, "Task completed and VERIFIED (deterministic acceptance).")
                break
            # N17 Phase B: the AGENT failed (e.g. the brain went down mid-way): say so with its REAL cause.
            # Before, the error was DISCARDED and the anti-loop ended up saying "the same missing checks
            # repeat ()", a misleading diagnosis that blamed the model when the brain wasn't responding. A
            # TRANSIENT failure keeps its retry (signature by exception type, like the missing checks); if
            # the error REPEATS (brain really down), status "error" with the cause.
            if _err:
                err_signature = ("__agent_error__", _err.split(":", 1)[0])
                if err_signature in signatures:
                    final_state, summary = "error", f"The agent failed: {_err}"
                    break
                signatures.append(err_signature)
                emit({"type": "worker_log", "text": f"⚠ the agent failed ({_err.split(':', 1)[0]}); retrying"})
                continue
            missing = tuple(txt for c, txt in details if not c)
            if missing in signatures:                          # ANTI-LOOP: same missing checks as before -> no progress
                final_state, summary = "stopped", "No progress: the same missing checks repeat (" + "; ".join(missing[:2]) + ")."
                break
            signatures.append(missing)
            current_task = task + " - fix what's missing: " + "; ".join(missing)   # feedback for the next pass
            # (N19: the REAL MODULE injection happens ABOVE, on EACH pass, no longer only here on retry)
    if final_state == "stopped" and not summary:
        summary = "Reached the maximum number of passes without meeting acceptance."
    return final_state, summary, attempt


def _cfg_num(cfg, key, default):
    try:
        v = (cfg.get(key) or "").strip() if hasattr(cfg, "get") else ""
        return type(default)(v) if v else default
    except (ValueError, TypeError):
        return default


_TERMINAL_OK_STATES = ("done", "done_unverified")   # stages the user can approve


def run_by_stages(task, stages, project, emit, cfg, router, resolve=None, signal_fn=None):
    """PHASE C: runs the task STAGE BY STAGE with the user's approval. Runs ONE stage (with the verified
    loop), STOPS, shows the result and waits for "approve / redo / cancel". Only continues on approval.
    Progress PERSISTS (.worker/stages.json) -> if it's interrupted, resending the SAME task RESUMES from the
    current stage (already-approved stages are not repeated). Reuses pause/resume."""
    import projects  # noqa: F401: consistency with the rest of the module (project_dir/base_dir)
    if resolve is None:
        resolve = _default_resolver
    project = (project or "").strip()
    if not project or project == "undefined":
        emit({"type": "worker_final", "state": "error",
              "summary": "Choose or create a project first (I didn't get a valid project)."})
        return
    if signal_fn is None:
        def signal_fn():
            return read_signal(cfg, project)

    # Per-batch stage CAP (light context): if the user asks for more, do the first MAX and WARN (at the top
    # and in the final summary), NEVER silently dropped (that would be a false pass by omission).
    max_stages = max(2, int(_cfg_num(cfg, "WORKER_MAX_STAGES", 8)))
    total_requested = len(stages)
    truncated = max(0, total_requested - max_stages)
    if truncated:
        stages = list(stages)[:max_stages]

    # RESUME after a crash: if there's persisted progress for THIS SAME task/stages, resume from the first
    # NOT-approved stage (approved ones left their artifact and aren't repeated).
    prev = read_stages(cfg, project)
    if prev.get("task") == task and prev.get("stages") == list(stages):
        states = (prev.get("estados") or ["pending"] * len(stages))[:len(stages)]
        states += ["pending"] * (len(stages) - len(states))
    else:
        states = ["pending"] * len(stages)
    i = next((k for k, e in enumerate(states) if e != "approved"), None)
    if i is None:                          # all approved = a COMPLETE previous run (not an interruption):
        states = ["pending"] * len(stages)   # resending the SAME task starts from SCRATCH, not "done" on the fly
        i = 0
    resuming = i > 0

    emit({"type": "worker_start", "task": task, "project": project})
    set_signal(cfg, project, "continue")
    # We persist the STAGES context already in the 'running' state (not only at the end): so if the user
    # RELOADS the view mid staged-run, restoreRun takes the by-stages branch and paints the list.
    save_state(cfg, project, state="running", task=task, attempts=0,
               plan=list(stages), steps=list(states), staged=True)
    emit({"type": "worker_stages", "stages": list(stages), "estados": states, "actual": i,
          "resuming": resuming})
    if truncated:
        emit({"type": "worker_log", "text": f"⚠️ you asked for {total_requested} stages; I'll do the first "
              f"{max_stages} this batch and let you know - the remaining {truncated} will NOT be done here "
              f"(send them in another batch)."})
    if resuming:
        emit({"type": "worker_log", "text": f"↩️ resuming from stage {i + 1}/{len(stages)} "
                                             f"(the {i} before it were already approved)."})
    proj_dir = project_dir(cfg, project)
    t0 = time.time()
    # N19: start instant of the RUN (not the stage). Contract-first uses it to re-read modules created in
    # ANY stage of this run, so the test stage DOES find the module from the previous stage (the N18 gap
    # between stages: the per-stage mtime filter discarded it).
    run_mtime = time.time() - 2

    while i < len(stages):
        if _wait_if_paused(signal_fn, emit) == "cancel":
            _finish_stages(cfg, project, emit, "stopped", "Cancelled by the user.", task, stages, states, i, t0)
            return
        states[i] = "in_progress"
        save_stages(cfg, project, task, stages, states, i)
        emit({"type": "worker_stage", "i": i, "state": "in_progress"})
        emit({"type": "worker_log", "text": f"━━ Stage {i + 1}/{len(stages)}: {stages[i][:90]} ━━"})
        cls = classify_task(stages[i])
        base_mtime = time.time() - 2               # the VERIFIER credits only what's created in THIS stage
        verifier = _verifier_for(cls, cfg, stages[i], since=base_mtime)
        state, summary, _attempt = _run_verified(
            stages[i], project, emit, cfg, router, resolve, signal_fn, cls, verifier, proj_dir,
            since=base_mtime, contract_since=run_mtime)   # contract-first re-reads modules from the WHOLE run
        states[i] = state
        save_stages(cfg, project, task, stages, states, i, summary)
        emit({"type": "worker_stage", "i": i, "state": state, "summary": summary})

        # STOP and wait for approval: the user is the checkpoint between stages. A stage that did NOT go well
        # (stopped/error) CANNOT be approved (approving a failure would give a false green "done"): only
        # redo/cancel are offered. Approval only makes sense on a terminal-OK state.
        can_approve = state in _TERMINAL_OK_STATES
        actions = ["approve", "redo", "cancel"] if can_approve else ["redo", "cancel"]
        emit({"type": "worker_awaiting_approval", "i": i, "state": state, "summary": summary,
              "last": i == len(stages) - 1, "actions": actions, "can_approve": can_approve})
        decision = _wait_for_approval(cfg, project, signal_fn, emit)
        if decision == "cancel":
            _finish_stages(cfg, project, emit, "stopped", f"Cancelled after stage {i + 1}.",
                           task, stages, states, i, t0)
            return
        # Defense in depth: an 'approve' signal on a NON-OK stage (tampered frontend / forced signal) is
        # treated as redo, a failure is never approved.
        if decision == "redo" or (decision == "approve" and not can_approve):
            if decision == "approve":
                emit({"type": "worker_log", "text": f"⚠️ stage {i + 1} ({state}) isn't ready to approve; redoing it"})
            else:
                emit({"type": "worker_log", "text": f"↻ redoing stage {i + 1}"})
            states[i] = "pending"
            save_stages(cfg, project, task, stages, states, i)
            continue                                    # re-run the SAME stage
        # approve (only possible on a terminal-OK state) -> mark as approved and move to the next
        states[i] = "approved"
        emit({"type": "worker_stage", "i": i, "state": "approved"})
        save_stages(cfg, project, task, stages, states, i + 1)
        i += 1

    if truncated:                                          # honest: not everything requested is done
        _finish_stages(cfg, project, emit, "done_unverified",
                       f"I did and you approved the first {len(stages)} stages, but you asked for "
                       f"{total_requested}: {truncated} are missing (send them in another batch).",
                       task, stages, states, len(stages), t0)
    else:
        _finish_stages(cfg, project, emit, "done",
                       f"The {len(stages)} stages completed and approved by you.", task, stages, states,
                       len(stages), t0)


def _finish_stages(cfg, project, emit, state, summary, task, stages, states, current, t0):
    """Common close of a staged run: emits worker_final, persists state.json (for the view's restoration)
    + stages.json, and writes the run log."""
    emit({"type": "worker_final", "state": state, "summary": summary})
    save_state(cfg, project, state=state, summary=summary, task=task[:300],
               seg=round(time.time() - t0, 1),
               plan=list(stages), steps=list(states), staged=True)
    save_stages(cfg, project, task, stages, states, current, summary)
    run_log(cfg, project,
            f"staged task '{task[:50]}' -> {state} · {summary} ({len(stages)} stages, {round(time.time()-t0)}s)")


def run(task, project, emit, cfg=None, router=None):
    """Worker entry point. With `router`: if the task carries STAGE markers (→ / numbered / "then") ->
    STAGED engine with approval (Phase C); otherwise -> single-task autonomous engine (Phase 3). Without
    `router` (UI tests) -> STUB engine."""
    if router is None:
        run_stub(task, project, emit)
        return
    # Reversibility (like WEB_SEARCH in Phase A): WORKER_STAGES="0" in config.env disables staged execution
    # -> every task goes to the usual autonomous engine (Phase 3), pre-Phase-C behavior.
    stages_on = str((cfg or {}).get("WORKER_STAGES", "1")).strip().lower() not in ("0", "false", "no", "off")
    stages = _split_into_stages(task) if stages_on else [task]
    if len(stages) >= 2:
        run_by_stages(task, stages, project, emit, cfg, router)
    else:
        run_autonomous(task, project, emit, cfg, router)
